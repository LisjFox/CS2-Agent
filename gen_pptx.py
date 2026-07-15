"""
生成 CS2 AI 教练展示 PPTX（基于教师要求大纲）
=========================================
5 个必选章节：
  1. 选题意义
  2. 方案设计
  3. 特色与创新
  4. 结果展示（含截图占位标注）
  5. 分工
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色 ──
BG_DARK    = RGBColor(0x0A, 0x0A, 0x0A)
BG_CARD    = RGBColor(0x14, 0x14, 0x14)
BLUE       = RGBColor(0x3B, 0x82, 0xF6)
PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_DIM   = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xD1, 0xD5, 0xDB)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_color(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//a:srgbClr', nsmap)
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_el.set('val', str(int(alpha * 1000)))
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=12, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '微软雅黑'
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bullet(tf, text, sub='', font_size=11, color=GRAY_LIGHT, sub_size=9):
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"▸ {text}"
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = True
    run.font.name = '微软雅黑'
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = f"   {sub}"
        r2.font.size = Pt(sub_size)
        r2.font.color.rgb = GRAY
        r2.font.name = '微软雅黑'
        p2.space_before = Pt(1)
        p.space_before = Pt(4)


def add_placeholder_box(slide, left, top, width, height, label, icon="📄"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # icon
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = icon
    r.font.size = Pt(28)
    r.font.color.rgb = GRAY
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = GRAY_DIM
    r2.font.name = '微软雅黑'
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "[截图待补充]"
    r3.font.size = Pt(14)
    r3.font.color.rgb = AMBER
    r3.font.bold = True
    r3.font.name = '微软雅黑'
    # center vertically
    tf.paragraphs[0].space_before = Pt(12)
    return shape


def chapter_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    # accent line
    add_shape_with_color(slide, Inches(1.5), Inches(2.6), Inches(6.5), Pt(3), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                f"0{number}" if number < 10 else str(number),
                font_size=12, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(8.5), Inches(1.2),
                title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(3.8), Inches(8.5), Inches(0.6),
                    subtitle, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # ════════════════════════════════════════════
    # Slide 1: 封面
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "CS2 · AI COACH", font_size=11, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5),
                "尚鸣惊人", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_shape_with_color(slide, Inches(5.5), Inches(4.0), Inches(2.3), Pt(3), BLUE)
    add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.8),
                "一站式 CS2 数据分析 + 战术教练系统",
                font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                "李尚基 × 潘一鸣    清华大学 · 2025",
                font_size=14, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # Slide 2: 选题意义 — 标题页
    # ════════════════════════════════════════════
    chapter_slide(prs, 1, "让 AI 当你的教练",
                  "CS2 每天上百万玩家打排位，但复盘？看几小时录像太累了")

    # ════════════════════════════════════════════
    # Slide 3: 为什么需要一个智能体
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "01 · 选题意义", font_size=11, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
                "为什么需要一个智能体？", font_size=28, bold=True, color=WHITE)
    items = [
        ("自动复盘", "3 秒解析整场 Demo，告别几小时录像回放"),
        ("个性化指导", "大模型根据你的数据生成针对性训练建议"),
        ("职业对标", "与 20+ 职业选手数据对比，找到差距在哪里"),
        ("实时指导", "GSI 实时监控，比赛中同步数据随时分析"),
    ]
    for i, (title, desc) in enumerate(items):
        y = Inches(2.2) + Inches(i * 1.1)
        add_shape_with_color(slide, Inches(1.2), y, Pt(4), Inches(0.8), BLUE)
        add_textbox(slide, Inches(1.6), y, Inches(3), Inches(0.4),
                    title, font_size=18, bold=True, color=WHITE)
        add_textbox(slide, Inches(1.6), y + Inches(0.35), Inches(9), Inches(0.5),
                    desc, font_size=12, color=GRAY)

    # ════════════════════════════════════════════
    # Slide 4: 方案设计 — 标题页
    # ════════════════════════════════════════════
    chapter_slide(prs, 2, "两个大脑，一个目标",
                  "输入：.dem 文件 / GSI 实时数据  →  输出：HTML 报告 + AI 对话 + 战术指导")

    # ════════════════════════════════════════════
    # Slide 5: 整体流程（步骤说明）
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "02 · 方案设计", font_size=11, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                "整体流程", font_size=28, bold=True, color=WHITE)
    steps = [
        ("01", "解析 Demo", "parser.py", BLUE),
        ("02", "分析 + 可视化", "analysis.py + visualization.py", PURPLE),
        ("03", "生成 HTML 报告", "report.py（图表全部内嵌，离线可用）", CYAN),
        ("04", "AI 对话 + 战术分析", "dialog.py + pan_coach.py（11 个 @tool + LLM）", AMBER),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        y = Inches(2.0) + Inches(i * 1.2)
        add_shape_with_color(slide, Inches(1.0), y, Pt(3), Inches(0.9), color)
        add_textbox(slide, Inches(0.5), y, Inches(0.8), Inches(0.9),
                    num, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.6), y, Inches(4), Inches(0.5),
                    title, font_size=18, bold=True, color=WHITE)
        add_textbox(slide, Inches(1.6), y + Inches(0.45), Inches(9), Inches(0.4),
                    desc, font_size=11, color=GRAY)

    # ════════════════════════════════════════════
    # Slide 6: 双 Agent 分栏
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "02 · 方案设计", font_size=11, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                "双 Agent 协作架构", font_size=28, bold=True, color=WHITE)

    col_left = Inches(2.0)
    col_right = Inches(7.0)
    col_w = Inches(4.5)
    h = Inches(5.0)
    y0 = Inches(1.8)

    # 左：李尚基
    add_shape_with_color(slide, col_left, y0, col_w, h, RGBColor(0x11, 0x1A, 0x2E))
    add_shape_with_color(slide, col_left, y0, col_w, Pt(3), BLUE)
    add_textbox(slide, col_left + Inches(0.3), y0 + Inches(0.2), col_w - Inches(0.6), Inches(0.6),
                "李尚基 Agent", font_size=20, bold=True, color=BLUE)
    add_textbox(slide, col_left + Inches(0.3), y0 + Inches(0.7), col_w - Inches(0.6), Inches(0.4),
                "数据分析引擎", font_size=11, color=RGBColor(0x93, 0xC5, 0xFD))
    li_items = [
        "Demo 解析（击杀/位置/武器/分段）",
        "数据可视化（Matplotlib 深色主题）",
        "HTML 报告生成",
        "长期 + 短期记忆系统",
        "20+ 职业选手 × 12 维对比",
        "GSI 实时监控",
    ]
    tf = add_textbox(slide, col_left + Inches(0.3), y0 + Inches(1.3), col_w - Inches(0.6), Inches(3.5),
                     "", font_size=11, color=GRAY_LIGHT)
    for item in li_items:
        add_para(tf, f"▸ {item}", font_size=11, color=GRAY_LIGHT)

    # 右：潘一鸣
    add_shape_with_color(slide, col_right, y0, col_w, h, RGBColor(0x1E, 0x12, 0x2E))
    add_shape_with_color(slide, col_right, y0, col_w, Pt(3), PURPLE)
    add_textbox(slide, col_right + Inches(0.3), y0 + Inches(0.2), col_w - Inches(0.6), Inches(0.6),
                "潘一鸣 Agent", font_size=20, bold=True, color=PURPLE)
    add_textbox(slide, col_right + Inches(0.3), y0 + Inches(0.7), col_w - Inches(0.6), Inches(0.4),
                "战术知识引擎", font_size=11, color=RGBColor(0xC4, 0xB5, 0xFD))
    pan_items = [
        "地图知识（6 张比赛地图点位）",
        "投掷物路线数据",
        "武器属性系统",
        "经济策略模型",
        "ChromaDB 向量检索（语义搜索）",
        "残局战术推理 + LLM 分析",
    ]
    tf = add_textbox(slide, col_right + Inches(0.3), y0 + Inches(1.3), col_w - Inches(0.6), Inches(3.5),
                     "", font_size=11, color=GRAY_LIGHT)
    for item in pan_items:
        add_para(tf, f"▸ {item}", font_size=11, color=GRAY_LIGHT)

    # 底部集成层说明
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "集成层 integration/ → IntegratedCoach 双 Agent 无缝配合",
                font_size=11, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # Slide 7: 11 个 @tool 工具
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "02 · 方案设计", font_size=11, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                "11 个 @tool 工具", font_size=28, bold=True, color=WHITE)

    li_tools = [
        ("🔍 get_player_history", "查询玩家历史记录", BLUE),
        ("📈 get_player_trend", "指标趋势分析", BLUE),
        ("👥 get_all_players", "查看所有玩家", BLUE),
        ("📋 get_training_plan", "每日训练计划", BLUE),
        ("🔫 get_weapon_advice", "武器训练建议", BLUE),
        ("🏆 compare_with_pro", "职业选手对比", BLUE),
        ("📊 get_pro_leaderboard", "职业选手排行榜", BLUE),
        ("📚 query_knowledge", "语义检索知识库", PURPLE),
        ("🗺️ get_map_tactics", "地图战术建议", PURPLE),
        ("💣 get_grenade_routes", "投掷物路线查询", PURPLE),
        ("💰 get_eco_advice", "经济策略分析", PURPLE),
    ]
    # 2-column grid
    col_w2 = Inches(5.0)
    for i, (name, desc, color) in enumerate(li_tools):
        col_idx = i // 6
        row_idx = i % 6
        x = Inches(1.2) + col_idx * Inches(5.8)
        y = Inches(1.8) + row_idx * Inches(0.8)
        bg = RGBColor(0x11, 0x1A, 0x2E) if color == BLUE else RGBColor(0x1E, 0x12, 0x2E)
        add_shape_with_color(slide, x, y, col_w2, Inches(0.65), bg)
        add_shape_with_color(slide, x, y, Pt(3), Inches(0.65), color)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.8), Inches(0.55),
                    name, font_size=11, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(2.9), y + Inches(0.08), Inches(2.1), Inches(0.5),
                    desc, font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "李尚基 × 7 个   |   潘一鸣 × 4 个   |   LangChain @tool 自动路由",
                font_size=11, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # Slide 8: 特色创新 — 标题页
    # ════════════════════════════════════════════
    chapter_slide(prs, 3, "七个杀手锏",
                  "不只是做出来了，还要做得漂亮")

    # ════════════════════════════════════════════
    # Slide 9: 20+ 职业选手 × 12 维
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "03 · 特色与创新", font_size=11, color=GRAY_DIM)

    add_textbox(slide, Inches(1.0), Inches(1.0), Inches(5), Inches(1.0),
                "20+", font_size=72, bold=True, color=WHITE)
    add_textbox(slide, Inches(3.5), Inches(1.2), Inches(3), Inches(0.6),
                "职业选手", font_size=16, color=GRAY)
    add_textbox(slide, Inches(3.5), Inches(1.7), Inches(3), Inches(0.5),
                "× 12 个维度", font_size=14, color=GRAY_DIM)

    dims = [
        ("K/D", "击杀比"), ("HS%", "爆头率"), ("ADR", "场均伤害"),
        ("Rating", "综合评分"), ("Impact", "影响力"), ("KPR", "回合击杀"),
        ("KAST", "存活贡献"), ("Clutch%", "残局胜率"), ("Multi%", "多杀率"),
        ("OpenKPR", "首杀率"), ("UtilDmg", "道具伤害"), ("Flash", "闪助攻"),
    ]
    for i, (abbr, name) in enumerate(dims):
        col = i % 6
        row = i // 6
        x = Inches(1.2) + col * Inches(1.9)
        y = Inches(3.5) + row * Inches(1.5)
        colors_cycle = [BLUE, PURPLE, CYAN, AMBER, GREEN, RGBColor(0xF4, 0x72, 0xB6)]
        c = colors_cycle[i % len(colors_cycle)]
        add_shape_with_color(slide, x, y, Inches(1.7), Inches(1.1), RGBColor(0x14, 0x14, 0x14))
        add_textbox(slide, x, y + Inches(0.1), Inches(1.7), Inches(0.5),
                    abbr, font_size=16, bold=True, color=c, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.6), Inches(1.7), Inches(0.4),
                    name, font_size=9, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "欧几里得距离 → 自动匹配与你打法最相似的选手",
                font_size=11, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # Slide 10: 七大创新点
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "03 · 特色与创新", font_size=11, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                "七大创新亮点", font_size=28, bold=True, color=WHITE)
    innovations = [
        ("双 Agent 协作架构", "数据分析 + 战术知识各司其职，通过集成层无缝配合", BLUE),
        ("职业选手对比系统", "20+ 选手 × 12 维度，欧几里得距离自动匹配相似选手", PURPLE),
        ("ChromaDB 向量知识库", "语义检索，支持按地图/类别过滤，90+ 知识条目", CYAN),
        ("GSI 实时监控", "比赛中实时推送数据，即时分析对局状态", AMBER),
        ("全自动 Demo 定位", "读 Steam 注册表 + libraryfolders.vdf，一键定位", GREEN),
        ("单文件 HTML 报告", "图表全部 base64 内嵌，离线可看，无网络依赖", RGBColor(0xF4, 0x72, 0xB6)),
        ("长短期双记忆", "JSON 长期记忆 + 对话上下文短期记忆，记录进步历程", RGBColor(0x81, 0x84, 0xF7)),
    ]
    for i, (title, desc, color) in enumerate(innovations):
        y = Inches(1.8) + Inches(i * 0.78)
        add_shape_with_color(slide, Inches(0.8), y, Pt(4), Inches(0.6), color)
        add_textbox(slide, Inches(1.3), y, Inches(3.5), Inches(0.35),
                    f"0{i+1}  {title}", font_size=15, bold=True, color=WHITE)
        add_textbox(slide, Inches(4.5), y + Inches(0.02), Inches(7.5), Inches(0.55),
                    desc, font_size=11, color=GRAY)

    # ════════════════════════════════════════════
    # Slide 11: 结果展示 — 标题页
    # ════════════════════════════════════════════
    chapter_slide(prs, 4, "用数据说话",
                  "所见即所得  ·  推荐准备视频备份以防现场状况")

    # ════════════════════════════════════════════
    # Slide 12-15: 截图占位
    # ════════════════════════════════════════════
    placeholders = [
        ("HTML 报告首页",
         "建议截取：报告总览页（K/D、HS%、ADR、Rating 等全景数据）",
         "📄"),
        ("击杀热力图 + AI 点评",
         "建议截取：Matplotlib 深色热力图 + LLM 生成的个性化分析文本",
         "🔥"),
        ("职业选手对比",
         "建议截取：对比详情 + 自动匹配的最相似 3 名选手排名",
         "🏆"),
        ("AI 对话 + GSI 实时监控",
         "建议截取：AI 教练对话界面（含 @tool 调用） + 实时数据推送",
         "💬"),
    ]
    for i, (title_hint, hint_text, icon) in enumerate(placeholders):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                    "04 · 结果展示", font_size=11, color=GRAY_DIM)
        add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                    title_hint, font_size=24, bold=True, color=WHITE)
        add_placeholder_box(slide, Inches(3.5), Inches(2.0), Inches(6.3), Inches(4.0),
                            hint_text, icon=icon)
        if i == 3:
            add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                        "11 个 @tool + LangChain 自动路由  |  GSI 实时推送血量/武器/位置",
                        font_size=10, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # Slide 16: 分工 — 标题页
    # ════════════════════════════════════════════
    chapter_slide(prs, 5, "各司其职",
                  "两人项目，各有所长")

    # ════════════════════════════════════════════
    # Slide 17: 双人分工详情
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                "05 · 分工", font_size=11, color=GRAY_DIM)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                "双人分工", font_size=28, bold=True, color=WHITE)

    col_left = Inches(1.5)
    col_right = Inches(7.0)
    col_w = Inches(4.5)
    y0 = Inches(1.8)
    box_h = Inches(5.0)

    # 左：李尚基
    add_shape_with_color(slide, col_left, y0, col_w, box_h, RGBColor(0x11, 0x1A, 0x2E))
    add_shape_with_color(slide, col_left, y0, col_w, Pt(3), BLUE)
    add_textbox(slide, col_left + Inches(0.3), y0 + Inches(0.2), col_w - Inches(0.6), Inches(0.6),
                "李尚基", font_size=22, bold=True, color=BLUE)
    add_textbox(slide, col_left + Inches(0.3), y0 + Inches(0.7), col_w - Inches(0.6), Inches(0.3),
                "学号：2025013285", font_size=10, color=GRAY_DIM)
    li_ls = [
        "Demo 解析引擎（parser.py）",
        "数据可视化（visualization.py）",
        "HTML 报告生成（report.py）",
        "LLM 对话系统（dialog.py）",
        "职业选手对比（pro_data.py）",
        "GSI 实时监控（live.py）",
        "记忆系统（memory.py）",
    ]
    tf = add_textbox(slide, col_left + Inches(0.3), y0 + Inches(1.2), col_w - Inches(0.6), Inches(3.5),
                     "", font_size=12, color=GRAY_LIGHT)
    for item in li_ls:
        add_para(tf, f"▸ {item}", font_size=12, color=GRAY_LIGHT, space_before=Pt(3))

    # 右：潘一鸣
    add_shape_with_color(slide, col_right, y0, col_w, box_h, RGBColor(0x1E, 0x12, 0x2E))
    add_shape_with_color(slide, col_right, y0, col_w, Pt(3), PURPLE)
    add_textbox(slide, col_right + Inches(0.3), y0 + Inches(0.2), col_w - Inches(0.6), Inches(0.6),
                "潘一鸣", font_size=22, bold=True, color=PURPLE)
    add_textbox(slide, col_right + Inches(0.3), y0 + Inches(0.7), col_w - Inches(0.6), Inches(0.3),
                "战术知识库 & 推理引擎", font_size=10, color=GRAY_DIM)
    li_pan = [
        "战术知识库（地图/投掷物/武器/经济）",
        "ChromaDB 向量检索（semantic search）",
        "战术推理引擎（pan_coach.py）",
        "工具函数模块",
        "地图点位分析",
        "经济策略 & 起枪建议",
    ]
    tf = add_textbox(slide, col_right + Inches(0.3), y0 + Inches(1.2), col_w - Inches(0.6), Inches(3.5),
                     "", font_size=12, color=GRAY_LIGHT)
    for item in li_pan:
        add_para(tf, f"▸ {item}", font_size=12, color=GRAY_LIGHT, space_before=Pt(3))

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "数据驱动 × 知识驱动 = 完整的 AI 教练",
                font_size=12, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # Slide 18: 结尾
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.4),
                "尚鸣惊人", font_size=11, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
                "你离职业\n只差一个 AI", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_shape_with_color(slide, Inches(5.5), Inches(4.0), Inches(2.3), Pt(3), BLUE)
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
                "李尚基 · 2025013285    潘一鸣",
                font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4),
                "清华大学 · 2025",
                font_size=13, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
                "谢谢！欢迎提问 🦊",
                font_size=14, color=GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # ── 保存 ──
    out_path = os.path.join(os.path.dirname(__file__), "尚鸣惊人_CS2_AI教练_展示.pptx")
    prs.save(out_path)
    print(f"[OK] PPTX 已生成: {out_path}")
    print(f"   共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    main()
